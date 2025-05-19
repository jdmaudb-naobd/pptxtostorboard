"""
PowerPoint (PPTX) Text and Image Extractor

This module provides functionality to extract text and images from PowerPoint (PPTX) files.
It uses python-pptx library for processing PowerPoint files and Pillow for image handling.
"""

import os
from typing import Dict, List, Tuple
from pptx import Presentation
from PIL import Image
from io import BytesIO

class PPTXExtractor:
    def __init__(self, pptx_path: str):
        """
        Initialize the PPTXExtractor with a PowerPoint file path.
        
        Args:
            pptx_path (str): Path to the PowerPoint file
        """
        if not os.path.exists(pptx_path):
            raise FileNotFoundError(f"PowerPoint file not found: {pptx_path}")
        
        self.pptx_path = pptx_path
        self.presentation = Presentation(pptx_path)

    def extract_text(self) -> List[Dict[str, str]]:
        """
        Extract text from all slides in the presentation.
        
        Returns:
            List[Dict[str, str]]: List of dictionaries containing slide number and text
        """
        text_content = []
        
        for idx, slide in enumerate(self.presentation.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
            
            text_content.append({
                "slide_number": idx,
                "text": "\n".join(slide_text)
            })
        
        return text_content

    def extract_images(self, output_dir: str) -> List[Dict[str, str]]:
        """
        Extract images from all slides and save them to the specified directory.
        
        Args:
            output_dir (str): Directory where images will be saved
            
        Returns:
            List[Dict[str, str]]: List of dictionaries containing image information
        """
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)
            
        image_info = []
        image_count = 0
        
        for slide_number, slide in enumerate(self.presentation.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    image_count += 1
                    image_bytes = shape.image.blob
                    image_type = shape.image.content_type.split('/')[-1]
                    image_filename = f"slide_{slide_number}_image_{image_count}.{image_type}"
                    image_path = os.path.join(output_dir, image_filename)
                    
                    # Save the image
                    with open(image_path, 'wb') as img_file:
                        img_file.write(image_bytes)
                    
                    # Get image dimensions
                    with Image.open(BytesIO(image_bytes)) as img:
                        width, height = img.size
                    
                    image_info.append({
                        "slide_number": slide_number,
                        "filename": image_filename,
                        "path": image_path,
                        "dimensions": f"{width}x{height}",
                        "format": image_type
                    })
        
        return image_info

def extract_all(pptx_path: str, images_output_dir: str) -> Tuple[List[Dict[str, str]], List[Dict[str, str]]]:
    """
    Convenience function to extract both text and images from a PowerPoint file.
    
    Args:
        pptx_path (str): Path to the PowerPoint file
        images_output_dir (str): Directory where images will be saved
        
    Returns:
        Tuple[List[Dict[str, str]], List[Dict[str, str]]]: Tuple containing text content and image information
    """
    extractor = PPTXExtractor(pptx_path)
    text_content = extractor.extract_text()
    image_info = extractor.extract_images(images_output_dir)
    return text_content, image_info 